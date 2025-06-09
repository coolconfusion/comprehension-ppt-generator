import os
import json
from pptx import Presentation
from pptx.util import Pt

def load_questions(json_path):
    """Load the JSON file of questions (list of dicts)."""
    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)

def create_ppt_split_by_category(template_path, questions, output_path):
    """
    Given a PPT template and a list of {question, category} dicts,
    generate a new PPT where each category gets its own slide.
    """
    prs = Presentation(template_path)

    # Assuming slide_layouts[1] is “Title and Content”
    title_and_content_layout = prs.slide_layouts[1]

    # Fixed category order
    categories = ["Literal", "Inferential", "Analytical", "Evaluative", "Appreciative"]

    # Group questions by category
    questions_by_cat = {cat: [] for cat in categories}
    for item in questions:
        cat = item["category"]
        text = item["question"]
        if cat in questions_by_cat:
            questions_by_cat[cat].append(text)
        else:
            questions_by_cat.setdefault(cat, []).append(text)

    # For each category, add a slide and fill in bullets
    for cat in categories:
        qlist = questions_by_cat.get(cat, [])
        if not qlist:
            continue

        slide = prs.slides.add_slide(title_and_content_layout)
        slide.shapes.title.text = f"{cat} Questions"
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        for q in qlist:
            p = text_frame.add_paragraph()
            p.text = q
            p.level = 0
            p.font.size = Pt(18)

    prs.save(output_path)
    print(f"Saved PPT: {output_path}")

if __name__ == "__main__":
    base_dir = os.path.dirname(__file__)            # e.g., H:\comprehension-ppt-generator\src
    template = os.path.join(base_dir, "..", "templates", "ppt_template.pptx")
    questions_json = os.path.join(base_dir, "..", "books", "output_questions.json")
    output_ppt = os.path.join(base_dir, "..", "generated_by_category.pptx")

    questions = load_questions(questions_json)
    create_ppt_split_by_category(template, questions, output_ppt)
